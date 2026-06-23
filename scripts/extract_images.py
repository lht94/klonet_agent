from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageOps


IMAGE_NAMESPACE = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
REL_NAMESPACE = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"


@dataclass
class ExtractedImage:
    source_path: str
    source_type: str
    collection: str
    category: str
    priority: str
    image_index: int
    image_path: str
    context_path: str
    width_px: int | None
    height_px: int | None
    sha256: str
    location: str
    nearby_text: str
    ocr_status: str
    ocr_text: str
    sensitivity: str
    notes: str
    ocr_language: str = ""
    ocr_confidence: float | None = None


def iter_manifest(path: Path) -> Iterable[dict]:
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            line = line.strip()
            if line:
                yield json.loads(line)


def slugify(value: str, max_length: int = 110) -> str:
    value = value.replace("\\", "/")
    value = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", value)
    value = value.strip("._")
    return value[:max_length] or "document"


def image_info(data: bytes) -> tuple[int | None, int | None]:
    try:
        with Image.open(__import__("io").BytesIO(data)) as image:
            return image.size
    except Exception:
        return None, None


def find_tesseract(explicit: str | None = None) -> str | None:
    candidates = [
        Path(explicit) if explicit else None,
        Path(os.environ["TESSERACT_CMD"]) if os.environ.get("TESSERACT_CMD") else None,
        Path(shutil.which("tesseract")) if shutil.which("tesseract") else None,
        Path(os.environ["LOCALAPPDATA"]) / "Programs" / "Tesseract-OCR" / "tesseract.exe"
        if os.environ.get("LOCALAPPDATA")
        else None,
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
    ]
    for candidate in candidates:
        if candidate and candidate.exists():
            return str(candidate)
    return None


def prepare_ocr_image(path: Path) -> Image.Image:
    with Image.open(path) as source:
        image = source.convert("RGB")
    if image.width < 1200 and image.height < 1200:
        image = image.resize((image.width * 2, image.height * 2), Image.Resampling.LANCZOS)
    return ImageOps.autocontrast(ImageOps.grayscale(image))


def rebuild_ocr_text(data: dict) -> str:
    lines: list[str] = []
    current_key: tuple[int, int, int, int] | None = None
    current_words: list[str] = []
    for index, raw_text in enumerate(data.get("text", [])):
        text = str(raw_text).strip()
        key = (
            int(data["page_num"][index]),
            int(data["block_num"][index]),
            int(data["par_num"][index]),
            int(data["line_num"][index]),
        )
        if current_key is not None and key != current_key and current_words:
            lines.append(" ".join(current_words))
            current_words = []
        current_key = key
        if text:
            current_words.append(text)
    if current_words:
        lines.append(" ".join(current_words))
    return "\n".join(lines)


def run_ocr(path: Path, language: str) -> dict:
    import pytesseract
    from pytesseract import Output

    image = prepare_ocr_image(path)
    data = pytesseract.image_to_data(
        image,
        lang=language,
        config="--oem 3 --psm 6",
        output_type=Output.DICT,
        timeout=90,
    )
    confidences = []
    for raw_conf, raw_text in zip(data.get("conf", []), data.get("text", [])):
        try:
            confidence = float(raw_conf)
        except (TypeError, ValueError):
            continue
        if confidence >= 0 and str(raw_text).strip():
            confidences.append(confidence)
    text = rebuild_ocr_text(data)
    return {
        "status": "ok" if text.strip() else "empty",
        "text": text,
        "confidence": round(sum(confidences) / len(confidences), 2) if confidences else None,
        "language": language,
        "error": "",
    }


def apply_ocr(
    root: Path,
    images: list[ExtractedImage],
    tesseract_cmd: str,
    language: str,
    workers: int,
) -> dict[str, int]:
    import pytesseract

    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    unique_images: dict[str, ExtractedImage] = {}
    for image in images:
        unique_images.setdefault(image.sha256, image)

    cache: dict[str, dict] = {}
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {
            executor.submit(run_ocr, root / image.image_path, language): sha
            for sha, image in unique_images.items()
        }
        for completed, future in enumerate(as_completed(futures), start=1):
            sha = futures[future]
            try:
                cache[sha] = future.result()
            except Exception as exc:
                cache[sha] = {
                    "status": "error",
                    "text": "",
                    "confidence": None,
                    "language": language,
                    "error": str(exc),
                }
            if completed % 25 == 0 or completed == len(futures):
                print(f"OCR progress: {completed}/{len(futures)}", flush=True)

    counts: dict[str, int] = {}
    for image in images:
        result = cache[image.sha256]
        image.ocr_status = f"ocr_{result['status']}"
        image.ocr_text = result["text"]
        image.ocr_language = result["language"]
        image.ocr_confidence = result["confidence"]
        if result["error"]:
            image.notes += f" OCR error: {result['error']}"
        elif result["status"] == "ok":
            image.notes += " OCR text is machine-generated and must be checked before promotion to curated knowledge."
        else:
            image.notes += " OCR was attempted but no text was extracted."
        write_context(image)
        counts[image.ocr_status] = counts.get(image.ocr_status, 0) + 1
    counts["unique_hashes_processed"] = len(cache)
    return counts

def write_context(record: ExtractedImage) -> None:
    path = Path(record.context_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        "\n".join(
            [
                "---",
                f"source_path: {record.source_path}",
                f"source_type: {record.source_type}",
                f"collection: {record.collection}",
                f"category: {record.category}",
                f"priority: {record.priority}",
                f"image_index: {record.image_index}",
                f"image_path: {record.image_path}",
                f"location: {record.location}",
                f"width_px: {record.width_px}",
                f"height_px: {record.height_px}",
                f"sha256: {record.sha256}",
                f"sensitivity: {record.sensitivity}",
                f"ocr_status: {record.ocr_status}",
                f"ocr_language: {record.ocr_language}",
                f"ocr_confidence: {record.ocr_confidence}",
                "---",
                "",
                "# Extracted Image Context",
                "",
                "## Nearby Text",
                "",
                record.nearby_text.strip() or "- no nearby text extracted",
                "",
                "## OCR Text",
                "",
                record.ocr_text.strip() or "- no OCR text extracted",
                "",
                "## Notes",
                "",
                record.notes.strip() or "- none",
                "",
            ]
        ),
        encoding="utf-8",
    )


def save_image(
    data: bytes,
    output_dir: Path,
    source_slug: str,
    index: int,
    extension: str,
) -> tuple[Path, str, int | None, int | None]:
    output_dir.mkdir(parents=True, exist_ok=True)
    sha = hashlib.sha256(data).hexdigest()
    ext = extension.lower().lstrip(".") or "bin"
    image_path = output_dir / f"{source_slug}__image_{index:04d}.{ext}"
    image_path.write_bytes(data)
    width, height = image_info(data)
    return image_path, sha, width, height


def text_from_element(element) -> str:
    texts = []
    for node in element.iter():
        if node.tag.endswith("}t") and node.text:
            text = node.text.strip()
            if text:
                texts.append(text)
    return " ".join(texts)


def extract_docx_images(root: Path, record: dict, output_root: Path) -> list[ExtractedImage]:
    from docx import Document

    source = root / record["source_path"]
    doc = Document(source)
    rels = doc.part.rels
    sequence: list[tuple[str, str | list[str]]] = []

    for child in doc.element.body.iterchildren():
        text = text_from_element(child)
        if text:
            sequence.append(("text", text))
        rel_ids: list[str] = []
        for blip in child.iter(IMAGE_NAMESPACE):
            rel_id = blip.get(REL_NAMESPACE)
            if rel_id:
                rel_ids.append(rel_id)
        if rel_ids:
            sequence.append(("image", rel_ids))

    images: list[ExtractedImage] = []
    source_slug = slugify(record["source_path"])
    image_index = 0
    for seq_index, (kind, value) in enumerate(sequence):
        if kind != "image":
            continue
        rel_ids = value if isinstance(value, list) else [value]
        for rel_id in rel_ids:
            if rel_id not in rels:
                continue
            target = rels[rel_id].target_part
            data = target.blob
            extension = Path(target.partname).suffix or ".bin"
            image_index += 1
            image_path, sha, width, height = save_image(data, output_root / "assets", source_slug, image_index, extension)
            before = [item[1] for item in sequence[max(0, seq_index - 4) : seq_index] if item[0] == "text"]
            after = [item[1] for item in sequence[seq_index + 1 : seq_index + 5] if item[0] == "text"]
            nearby = "\n\n".join([str(item) for item in before + after])
            context_path = output_root / "contexts" / f"{source_slug}__image_{image_index:04d}.context.md"
            extracted = ExtractedImage(
                source_path=record["source_path"],
                source_type=record["type"],
                collection=record["collection"],
                category=record["category"],
                priority=record["priority"],
                image_index=image_index,
                image_path=rel_path(root, image_path),
                context_path=rel_path(root, context_path),
                width_px=width,
                height_px=height,
                sha256=sha,
                location=f"docx body image occurrence {image_index}",
                nearby_text=nearby,
                ocr_status="ocr_not_available",
                ocr_text="",
                sensitivity="review_required",
                notes="Extracted embedded DOCX image evidence; review both the image and OCR text before using it in public knowledge.",
            )
            write_context(extracted)
            images.append(extracted)
    return images


def rel_path(root: Path, path: Path) -> str:
    return str(path.resolve().relative_to(root.resolve())).replace("\\", "/")


def extract_pptx_images(root: Path, record: dict, output_root: Path) -> list[ExtractedImage]:
    from pptx import Presentation

    source = root / record["source_path"]
    prs = Presentation(source)
    images: list[ExtractedImage] = []
    source_slug = slugify(record["source_path"])
    image_index = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    slide_texts.append(text)
        nearby = "\n\n".join(slide_texts)
        for shape in slide.shapes:
            rel_ids: list[str] = []
            for blip in shape.element.iter(IMAGE_NAMESPACE):
                rel_id = blip.get(REL_NAMESPACE)
                if rel_id:
                    rel_ids.append(rel_id)
            if not rel_ids:
                continue
            for rel_id in rel_ids:
                if rel_id not in slide.part.rels:
                    continue
                target = slide.part.rels[rel_id].target_part
                data = target.blob
                extension = Path(target.partname).suffix or ".bin"
                image_index += 1
                image_path, sha, width, height = save_image(data, output_root / "assets", source_slug, image_index, extension)
                context_path = output_root / "contexts" / f"{source_slug}__image_{image_index:04d}.context.md"
                extracted = ExtractedImage(
                    source_path=record["source_path"],
                    source_type=record["type"],
                    collection=record["collection"],
                    category=record["category"],
                    priority=record["priority"],
                    image_index=image_index,
                    image_path=rel_path(root, image_path),
                    context_path=rel_path(root, context_path),
                    width_px=width,
                    height_px=height,
                    sha256=sha,
                    location=f"slide {slide_number}",
                    nearby_text=nearby,
                    ocr_status="ocr_not_available",
                    ocr_text="",
                    sensitivity="review_required",
                    notes="Extracted PPTX picture evidence; review both the image and OCR text before using it in public knowledge.",
                )
                write_context(extracted)
                images.append(extracted)
    return images


def extract_pdf_pages(root: Path, record: dict, output_root: Path) -> list[ExtractedImage]:
    source = root / record["source_path"]
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return []

    source_slug = slugify(record["source_path"])
    text_pages = {}
    try:
        import pdfplumber

        with pdfplumber.open(source) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                text_pages[page_number] = (page.extract_text() or "").strip()
    except Exception:
        text_pages = {}

    images: list[ExtractedImage] = []
    with tempfile.TemporaryDirectory() as temp_dir:
        prefix = str(Path(temp_dir) / "page")
        subprocess.run([pdftoppm, "-png", "-r", "150", str(source), prefix], check=True, capture_output=True)
        for page_file in sorted(Path(temp_dir).glob("page-*.png")):
            match = re.search(r"page-(\d+)\.png$", page_file.name)
            page_number = int(match.group(1)) if match else len(images) + 1
            data = page_file.read_bytes()
            image_index = len(images) + 1
            image_path, sha, width, height = save_image(data, output_root / "assets", source_slug, image_index, ".png")
            context_path = output_root / "contexts" / f"{source_slug}__page_{page_number:04d}.context.md"
            extracted = ExtractedImage(
                source_path=record["source_path"],
                source_type=record["type"],
                collection=record["collection"],
                category=record["category"],
                priority=record["priority"],
                image_index=image_index,
                image_path=rel_path(root, image_path),
                context_path=rel_path(root, context_path),
                width_px=width,
                height_px=height,
                sha256=sha,
                location=f"pdf page {page_number}",
                nearby_text=text_pages.get(page_number, ""),
                ocr_status="ocr_not_available",
                ocr_text="",
                sensitivity="review_required",
                notes="Rendered PDF page image with pdftoppm; extracted page text is stored as nearby text when available.",
            )
            write_context(extracted)
            images.append(extracted)
    return images


def extract_images(root: Path, record: dict, output_root: Path) -> list[ExtractedImage]:
    if record["type"] == "docx":
        return extract_docx_images(root, record, output_root)
    if record["type"] == "pptx":
        return extract_pptx_images(root, record, output_root)
    if record["type"] == "pdf":
        return extract_pdf_pages(root, record, output_root)
    return []


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract image evidence for priority raw knowledge documents.")
    parser.add_argument("--manifest", default="knowledge/raw_manifest.jsonl")
    parser.add_argument("--root", default=".")
    parser.add_argument("--output", default="knowledge/extracted_images")
    parser.add_argument("--priorities", nargs="+", default=["P0", "P1"])
    parser.add_argument("--tesseract-cmd")
    parser.add_argument("--ocr-lang", default="chi_sim+eng")
    parser.add_argument("--skip-ocr", action="store_true")
    parser.add_argument("--ocr-workers", type=int, default=4)
    args = parser.parse_args()

    root = Path(args.root).resolve()
    manifest = root / args.manifest
    output_root = root / args.output
    output_root.mkdir(parents=True, exist_ok=True)

    all_images: list[ExtractedImage] = []
    source_summaries: list[dict] = []
    for record in iter_manifest(manifest):
        if record.get("priority") not in args.priorities:
            continue
        if record.get("type") not in {"docx", "pptx", "pdf"}:
            continue
        try:
            images = extract_images(root, record, output_root)
            status = "ok"
            error = ""
        except Exception as exc:  # pragma: no cover - best-effort batch command
            images = []
            status = "error"
            error = str(exc)
        all_images.extend(images)
        source_summaries.append(
            {
                "source_path": record["source_path"],
                "type": record["type"],
                "priority": record["priority"],
                "status": status,
                "image_count": len(images),
                "error": error,
            }
        )

    tesseract_cmd = None if args.skip_ocr else find_tesseract(args.tesseract_cmd)
    if tesseract_cmd:
        ocr_counts = apply_ocr(root, all_images, tesseract_cmd, args.ocr_lang, args.ocr_workers)
    else:
        ocr_counts = {"ocr_not_available": len(all_images), "unique_hashes_processed": 0}

    manifest_path = output_root / "image_manifest.jsonl"
    with manifest_path.open("w", encoding="utf-8") as handle:
        for image in all_images:
            handle.write(json.dumps(image.__dict__, ensure_ascii=False, separators=(",", ":")) + "\n")

    summary_path = output_root / "_image_extraction_summary.json"
    summary_path.write_text(json.dumps(source_summaries, ensure_ascii=False, indent=2), encoding="utf-8")
    print(
        json.dumps(
            {
                "sources": len(source_summaries),
                "images": len(all_images),
                "manifest": str(manifest_path),
                "summary": str(summary_path),
                "ocr": ocr_counts,
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
