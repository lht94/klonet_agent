from __future__ import annotations

import argparse
import html.parser
import json
import os
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable


class PlainTextHTMLParser(html.parser.HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() in {"p", "br", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)

    def text(self) -> str:
        lines = []
        for raw in "\n".join(self.parts).splitlines():
            line = " ".join(raw.split())
            if line:
                lines.append(line)
        return "\n".join(lines)


@dataclass
class OCRContext:
    """OCR 运行上下文：把依赖检测结果显式传给各类文档抽取器。"""

    enabled: bool
    lang: str
    available: bool
    warnings: list[str]


def build_ocr_context(enabled: bool, lang: str) -> OCRContext:
    if not enabled:
        return OCRContext(enabled=False, lang=lang, available=False, warnings=[])

    warnings: list[str] = []
    try:
        import pytesseract

        configure_tesseract_cmd(pytesseract)
        pytesseract.get_tesseract_version()
    except Exception as exc:  # pragma: no cover - depends on local OCR tools
        warnings.append(f"OCR requested but Tesseract is not available: {exc}")
        return OCRContext(enabled=True, lang=lang, available=False, warnings=warnings)

    try:
        from PIL import Image  # noqa: F401
    except Exception as exc:  # pragma: no cover - depends on local packages
        warnings.append(f"OCR requested but Pillow is not available: {exc}")
        return OCRContext(enabled=True, lang=lang, available=False, warnings=warnings)

    return OCRContext(enabled=True, lang=lang, available=True, warnings=[])


def configure_tesseract_cmd(pytesseract_module: object) -> None:
    """允许 Windows 用户通过环境变量或默认安装路径找到 tesseract.exe。"""

    command = os.environ.get("TESSERACT_CMD")
    local_appdata = os.environ.get("LOCALAPPDATA")
    candidates = [
        Path(command) if command else None,
        Path(local_appdata) / "Programs" / "Tesseract-OCR" / "tesseract.exe" if local_appdata else None,
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
    ]
    for candidate in candidates:
        if candidate and candidate.exists():
            pytesseract_module.pytesseract.tesseract_cmd = str(candidate)
            return


def ocr_image_bytes(blob: bytes, context: OCRContext, label: str) -> tuple[str, list[str]]:
    if not context.enabled:
        return "", []
    if not context.available:
        return "", context.warnings.copy()

    warnings: list[str] = []
    try:
        import pytesseract
        from PIL import Image

        configure_tesseract_cmd(pytesseract)
        with Image.open(BytesIO(blob)) as image:
            text = pytesseract.image_to_string(image, lang=context.lang)
        return text.strip(), warnings
    except Exception as exc:  # pragma: no cover - depends on local OCR tools
        warnings.append(f"OCR failed for {label}: {exc}")
        return "", warnings


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
        return path.read_bytes().decode("utf-8", errors="replace")


def extract_docx(path: Path, ocr: OCRContext | None = None) -> tuple[str, list[str]]:
    from docx import Document

    doc = Document(path)
    chunks: list[str] = []
    warnings: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)

    for table_index, table in enumerate(doc.tables, start=1):
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([" ".join(cell.text.split()) for cell in row.cells])
        if not rows:
            continue
        chunks.append(f"\n### Table {table_index}\n")
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        chunks.append("| " + " | ".join(normalized[0]) + " |")
        chunks.append("| " + " | ".join(["---"] * width) + " |")
        for row in normalized[1:]:
            chunks.append("| " + " | ".join(row) + " |")

    rels = getattr(doc.part, "rels", {})
    image_rels = [rel for rel in rels.values() if "image" in rel.reltype]
    image_count = len(image_rels)
    if image_count:
        if ocr and ocr.enabled:
            ocr_chunks: list[str] = []
            for image_index, rel in enumerate(image_rels, start=1):
                blob = getattr(rel.target_part, "blob", b"")
                text, image_warnings = ocr_image_bytes(blob, ocr, f"{path.name} image {image_index}")
                warnings.extend(image_warnings)
                if text:
                    ocr_chunks.append(f"### OCR Image {image_index}\n\n{text}")
            if ocr_chunks:
                chunks.append("\n## Embedded Image OCR\n")
                chunks.extend(ocr_chunks)
            else:
                warnings.append(f"document contains {image_count} embedded image(s); no OCR text extracted")
        else:
            warnings.append(f"document contains {image_count} embedded image(s); image OCR evidence is stored separately under knowledge/extracted_images when available")

    return "\n\n".join(chunks), warnings


def extract_pdf(path: Path, ocr: OCRContext | None = None) -> tuple[str, list[str]]:
    warnings: list[str] = []
    try:
        import pdfplumber

        pages: list[str] = []
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"## Page {index}\n\n{text.strip()}")
        if pages:
            return "\n\n".join(pages), warnings
        warnings.append("pdfplumber returned no text; tried pypdf fallback")
    except Exception as exc:  # pragma: no cover - best-effort extractor
        warnings.append(f"pdfplumber failed: {exc}")

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"## Page {index}\n\n{text.strip()}")
        text = "\n\n".join(pages)
        if text.strip():
            return text, warnings
        if ocr and ocr.enabled:
            return extract_pdf_with_ocr(path, ocr, warnings)
        return text, warnings
    except Exception as exc:  # pragma: no cover - best-effort extractor
        warnings.append(f"pypdf failed: {exc}")
        if ocr and ocr.enabled:
            return extract_pdf_with_ocr(path, ocr, warnings)
        return "", warnings


def extract_pdf_with_ocr(path: Path, ocr: OCRContext, warnings: list[str]) -> tuple[str, list[str]]:
    if not ocr.available:
        warnings.extend(ocr.warnings)
        return "", warnings

    try:
        from pdf2image import convert_from_path
    except Exception as exc:  # pragma: no cover - depends on local packages
        warnings.append(f"PDF OCR requested but pdf2image is not available: {exc}")
        return "", warnings

    try:
        import pytesseract

        pages = convert_from_path(str(path), dpi=200)
        ocr_pages: list[str] = []
        for page_index, image in enumerate(pages, start=1):
            text = pytesseract.image_to_string(image, lang=ocr.lang).strip()
            if text:
                ocr_pages.append(f"## OCR Page {page_index}\n\n{text}")
        if not ocr_pages:
            warnings.append("PDF OCR finished but no text was extracted")
        return "\n\n".join(ocr_pages), warnings
    except Exception as exc:  # pragma: no cover - depends on local OCR tools
        warnings.append(f"PDF OCR failed, check Poppler and Tesseract installation: {exc}")
        return "", warnings


def extract_pptx(path: Path, ocr: OCRContext | None = None) -> tuple[str, list[str]]:
    from pptx import Presentation

    prs = Presentation(path)
    slides: list[str] = []
    warnings: list[str] = []
    image_count = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    texts.append(text)
            if hasattr(shape, "image"):
                image_count += 1
                if ocr and ocr.enabled:
                    text, image_warnings = ocr_image_bytes(
                        shape.image.blob,
                        ocr,
                        f"{path.name} slide {slide_index} image {image_count}",
                    )
                    warnings.extend(image_warnings)
                    if text:
                        texts.append(f"### OCR Image {image_count}\n\n{text}")
        if texts:
            slides.append(f"## Slide {slide_index}\n\n" + "\n\n".join(texts))

    if image_count:
        if ocr and ocr.enabled:
            if not any("### OCR Image" in slide for slide in slides):
                warnings.append(f"presentation contains {image_count} picture shape(s); no OCR text extracted")
        else:
            warnings.append(f"presentation contains {image_count} picture shape(s); image OCR evidence is stored separately under knowledge/extracted_images when available")
    return "\n\n".join(slides), warnings


def extract_html(path: Path) -> str:
    parser = PlainTextHTMLParser()
    parser.feed(read_text_file(path))
    return parser.text()


def extract_record(source: Path, doc_type: str, ocr: OCRContext | None = None) -> tuple[str, list[str]]:
    if doc_type in {"markdown", "text", "rst"}:
        return read_text_file(source), []
    if doc_type == "html":
        return extract_html(source), []
    if doc_type == "docx":
        return extract_docx(source, ocr)
    if doc_type == "pdf":
        return extract_pdf(source, ocr)
    if doc_type == "pptx":
        return extract_pptx(source, ocr)
    return "", [f"unsupported extractable type: {doc_type}"]


def iter_manifest(path: Path) -> Iterable[dict]:
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            line = line.strip()
            if line:
                yield json.loads(line)


def markdown_header(record: dict, warnings: list[str]) -> str:
    target_layers = ", ".join(record.get("target_layers") or [])
    warning_text = "\n".join(f"- {item}" for item in warnings) if warnings else "- none"
    return "\n".join(
        [
            "---",
            f"source_path: {record['source_path']}",
            f"collection: {record['collection']}",
            f"type: {record['type']}",
            f"category: {record['category']}",
            f"priority: {record['priority']}",
            f"quality: {record['quality']}",
            f"sensitivity: {record['sensitivity']}",
            f"target_layers: {target_layers}",
            f"size_bytes: {record['size_bytes']}",
            f"last_modified: {record['last_modified']}",
            "---",
            "",
            "# Extracted Raw Text",
            "",
            "## Extraction Notes",
            "",
            warning_text,
            "",
            "## Content",
            "",
        ]
    )


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract raw evidence files into Markdown text intermediates.")
    parser.add_argument("--manifest", default="knowledge/raw_manifest.jsonl")
    parser.add_argument("--root", default=".")
    parser.add_argument("--summary", default="knowledge/extracted_docs/_extraction_summary.json")
    parser.add_argument("--ocr", action="store_true", help="Enable best-effort OCR for embedded images and scanned PDFs.")
    parser.add_argument("--ocr-lang", default="chi_sim+eng", help="Tesseract language setting, for example chi_sim+eng.")
    args = parser.parse_args()

    root = Path(args.root).resolve()
    manifest = (root / args.manifest).resolve()
    summary_path = (root / args.summary).resolve()
    ocr = build_ocr_context(args.ocr, args.ocr_lang)

    results: list[dict] = []
    for record in iter_manifest(manifest):
        derived_text = record.get("derived_text")
        if not derived_text:
            continue

        source = (root / record["source_path"]).resolve()
        output = (root / derived_text).resolve()
        output.parent.mkdir(parents=True, exist_ok=True)

        status = "ok"
        warnings: list[str] = []
        content = ""
        try:
            content, warnings = extract_record(source, record["type"], ocr)
            if not content.strip():
                status = "empty"
                warnings.append("no extractable text found")
        except Exception as exc:  # pragma: no cover - best-effort batch command
            status = "error"
            warnings.append(str(exc))

        output.write_text(markdown_header(record, warnings) + content.strip() + "\n", encoding="utf-8")
        results.append(
            {
                "source_path": record["source_path"],
                "derived_text": derived_text,
                "type": record["type"],
                "priority": record["priority"],
                "status": status,
                "warnings": warnings,
            }
        )

    summary_path.parent.mkdir(parents=True, exist_ok=True)
    summary_path.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
    counts: dict[str, int] = {}
    for result in results:
        counts[result["status"]] = counts.get(result["status"], 0) + 1
    print(json.dumps({"total": len(results), "status_counts": counts, "summary": str(summary_path)}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
